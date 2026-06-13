"""Text extraction from uploaded files (PDF, Word, and plain-text/code formats)."""

import fitz  # PyMuPDF


def extract_text_from_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    try:
        return "".join(page.get_text() for page in doc)
    finally:
        doc.close()  # release the file handle (avoids a leak)


def extract_text_from_docx(file_path: str) -> str:
    # Imported lazily so a missing optional dependency never crashes startup —
    # only .docx uploads fail (with a clear message) if python-docx isn't installed.
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError(
            "Word (.docx) support needs python-docx. Install it with: pip install python-docx"
        )
    doc = Document(file_path)
    parts = [p.text for p in doc.paragraphs]
    # Include table cells too — resumes/reports often use tables.
    for table in doc.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def extract_text_from_xlsx(file_path: str) -> str:
    """Extract all cell text from an .xlsx workbook (lazy import — optional dep)."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise RuntimeError("Excel (.xlsx) support needs openpyxl.")
    wb = load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    try:
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    parts.append("\t".join(cells))
    finally:
        wb.close()
    return "\n".join(parts)


def extract_text_from_pptx(file_path: str) -> str:
    """Extract slide text (and table cells) from a .pptx deck (lazy import)."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("PowerPoint (.pptx) support needs python-pptx.")
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text for c in row.cells]
                    if any(cells):
                        parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_text_plain(file_path: str) -> str:
    with open(file_path, "rb") as f:
        raw = f.read()
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def extract_text_from_file(file_path: str, ext: str) -> str:
    """Extract text from a file based on its extension.
    .pdf -> PyMuPDF, .docx -> python-docx, everything else -> decode as text."""
    ext = (ext or "").lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    if ext == ".docx":
        return extract_text_from_docx(file_path)
    if ext == ".xlsx":
        return extract_text_from_xlsx(file_path)
    if ext == ".pptx":
        return extract_text_from_pptx(file_path)
    # .txt, .md, .csv, .json, and source-code files are read as plain text.
    return _extract_text_plain(file_path)
