"""Media + document generation:
- Images via NVIDIA NIM FLUX.1-schnell (same key as chat — no extra signup).
- Videos via Pollinations.ai (no signup, generation triggered by the GET URL).
- PDFs via xhtml2pdf (LLM writes Markdown -> we render to a styled PDF).

All outbound calls are blocking; route handlers wrap them in run_in_threadpool.
"""

import io
import urllib.parse

import requests

from app.core.config import NVIDIA_API_KEY, POLLINATIONS_API_KEY


# ── Image: NVIDIA NIM FLUX.1-schnell ─────────────────────────────────────────
_IMAGE_URL = "https://ai.api.nvidia.com/v1/genai/black-forest-labs/flux.1-schnell"
# NIM FLUX accepts only this discrete set of side lengths — anything else 422s.
_FLUX_ALLOWED_SIDES = (768, 832, 896, 960, 1024, 1088, 1152, 1216, 1280, 1344)


def _clamp_flux_side(v: int) -> int:
    """Snap a requested side length to the nearest NIM-FLUX-supported value."""
    v = int(v or 1024)
    return min(_FLUX_ALLOWED_SIDES, key=lambda x: abs(x - v))


def _mime_from_b64(b64: str) -> str:
    """Detect the actual image format from the first chars of base64 data.
    NVIDIA NIM FLUX returns JPEG bytes; SD models often return PNG. Always
    label the data: URI with the right MIME so browsers render reliably."""
    if not b64:
        return "image/png"
    # JPEG magic FF D8 FF → base64 starts with '/9j/' or rarely '/9k/'
    if b64.startswith("/9j/") or b64.startswith("/9k/"):
        return "image/jpeg"
    # PNG magic 89 50 4E 47 → base64 starts with 'iVBOR'
    if b64.startswith("iVBOR"):
        return "image/png"
    # GIF magic 47 49 46 38 → 'R0lGO'
    if b64.startswith("R0lGO"):
        return "image/gif"
    # WEBP magic 52 49 46 46 (RIFF) → 'UklGR'
    if b64.startswith("UklGR"):
        return "image/webp"
    return "image/png"  # safe default


def _extract_b64(payload):
    """The NIM response wraps the base64 image in one of several shapes; try them all."""
    if not isinstance(payload, dict):
        return None
    # Direct keys
    for key in ("image", "b64_json", "base64", "data"):
        v = payload.get(key)
        if isinstance(v, str) and len(v) > 100:
            return v
    # artifacts[].base64 / .b64_json (most common NIM shape)
    arts = payload.get("artifacts")
    if isinstance(arts, list) and arts:
        a = arts[0]
        if isinstance(a, dict):
            for k in ("base64", "b64_json", "image"):
                v = a.get(k)
                if isinstance(v, str) and len(v) > 100:
                    return v
    # data[].b64_json (OpenAI-style)
    data = payload.get("data")
    if isinstance(data, list) and data:
        d = data[0]
        if isinstance(d, dict):
            for k in ("b64_json", "base64", "image"):
                v = d.get(k)
                if isinstance(v, str) and len(v) > 100:
                    return v
    return None


def generate_image_b64(prompt: str, width: int = 1024, height: int = 1024) -> str:
    """Generate an image via NVIDIA NIM FLUX.1-schnell. Returns a data: URI."""
    if not NVIDIA_API_KEY:
        raise RuntimeError("NVIDIA_API_KEY is not configured.")

    # Snap requested dimensions to the model's supported discrete set.
    w = _clamp_flux_side(width)
    h = _clamp_flux_side(height)

    r = requests.post(
        _IMAGE_URL,
        headers={
            "Authorization": f"Bearer {NVIDIA_API_KEY}",
            "Accept": "application/json",
            "Content-Type": "application/json",
        },
        json={
            "prompt": prompt,
            "width": w,
            "height": h,
            "seed": 0,
            "steps": 4,
        },
        timeout=90,
    )
    if r.status_code >= 400:
        raise RuntimeError(f"NVIDIA NIM image API returned {r.status_code}: {r.text[:300]}")
    try:
        payload = r.json()
    except Exception:
        raise RuntimeError("NVIDIA NIM image API returned non-JSON response.")

    b64 = _extract_b64(payload)
    if not b64:
        raise RuntimeError("NVIDIA NIM image API returned no base64 image data.")
    # Strip any whitespace (MIME-style line breaks every 76 chars would break
    # downstream consumers like markdown URL parsing or strict base64 decoders).
    b64 = "".join(b64.split())
    mime = _mime_from_b64(b64)
    return f"data:{mime};base64,{b64}"


# ── Video: Pollinations.ai (free signup required as of mid-2026) ─────────────
def generate_video_bytes(
    prompt: str, model: str = "wan-fast", duration: int = 5
) -> tuple[bytes, str]:
    """Generate a video via Pollinations.ai. Returns (bytes, content_type).

    The Pollinations key is sent ONLY as a server-side Bearer header — never
    embedded in any URL the browser sees. We pay the cost of proxying the MP4
    bytes through our backend to keep the key fully secret.

    Sign up free at https://enter.pollinations.ai/ and set POLLINATIONS_API_KEY
    in .env. Seed tier (the free signup) is unmetered on free models like
    wan-fast / wan / ltx-2.
    """
    if not POLLINATIONS_API_KEY:
        raise RuntimeError(
            "Video generation needs a free Pollinations.ai key. "
            "Sign up at https://enter.pollinations.ai/ and set POLLINATIONS_API_KEY in .env."
        )
    encoded = urllib.parse.quote(prompt.strip())[:500]
    dur = max(2, min(int(duration or 5), 8))
    url = f"https://gen.pollinations.ai/video/{encoded}?model={model}&duration={dur}"

    r = requests.get(
        url,
        headers={"Authorization": f"Bearer {POLLINATIONS_API_KEY}"},
        timeout=180,
    )
    if r.status_code >= 400:
        raise RuntimeError(
            f"Pollinations video API returned {r.status_code}: {r.text[:300]}"
        )
    content_type = r.headers.get("Content-Type", "video/mp4")
    return r.content, content_type


# ── PDF: LLM -> Markdown -> styled HTML -> xhtml2pdf -> PDF ──────────────────
_PDF_CSS = """
@page { size: A4; margin: 1.8cm 2cm; }
body { font-family: Helvetica, Arial, sans-serif; line-height: 1.55; color: #1a1a1a; font-size: 11pt; }
h1 { color: #0a0a0a; border-bottom: 2px solid #444; padding-bottom: 6px; font-size: 22pt; margin: 0 0 14pt 0; }
h2 { color: #1a1a1a; font-size: 15pt; margin: 22pt 0 8pt 0; border-bottom: 1px solid #e2e2e2; padding-bottom: 4px; }
h3 { color: #2a2a2a; font-size: 12.5pt; margin: 16pt 0 6pt 0; }
p  { margin: 0 0 8pt 0; }
ul, ol { margin: 0 0 10pt 0; padding-left: 20pt; }
li { margin: 0 0 3pt 0; }
code { background: #f4f4f6; padding: 1px 4px; border-radius: 3px; font-family: Courier, monospace; font-size: 10pt; }
pre { background: #f4f4f6; padding: 10px 12px; border-radius: 4px; font-family: Courier, monospace; font-size: 9.5pt; white-space: pre-wrap; }
blockquote { border-left: 4px solid #ddd; margin: 0 0 10pt 0; padding: 2pt 12pt; color: #555; }
table { border-collapse: collapse; width: 100%; margin: 8pt 0 12pt 0; }
th, td { border: 1px solid #ddd; padding: 6pt 8pt; text-align: left; vertical-align: top; }
th { background: #f5f5f5; }
hr { border: 0; border-top: 1px solid #ddd; margin: 18pt 0; }
a { color: #2563eb; text-decoration: underline; }
"""


def _ensure_table_spacing(md: str) -> str:
    """Guarantee a blank line before and after every Markdown table block.

    The Python ``markdown`` tables extension (part of the ``extra`` bundle)
    requires a blank line *before* the first ``|`` row and another blank line
    *after* the last row.  LLMs often run a table directly under a heading
    without the required blank lines, which makes the parser skip the table
    and emit raw ``| … |`` pipe characters instead.  This one-pass fixer
    inserts the missing blank lines before we hand the text to the renderer.
    """
    lines = md.split('\n')
    out: list = []
    for i, line in enumerate(lines):
        is_table = line.strip().startswith('|')
        prev_table = bool(out) and out[-1].strip().startswith('|')

        # Need a blank line *before* the first row of a new table block.
        if is_table and not prev_table and out and out[-1].strip():
            out.append('')

        out.append(line)

        # Need a blank line *after* the last row when the next line is
        # non-empty content (not another table row).
        if is_table:
            nxt = lines[i + 1].strip() if i + 1 < len(lines) else ''
            if nxt and not nxt.startswith('|'):
                out.append('')

    return '\n'.join(out)


def _md_to_html(md_text: str) -> str:
    """Markdown → HTML using the ``markdown`` package (extra + fenced code)."""
    import markdown

    # Fix blank-line spacing around tables BEFORE the parser sees the text —
    # the tables extension inside `extra` silently skips tables that don't
    # have blank lines before/after them, rendering raw ``| … |`` instead.
    spaced = _ensure_table_spacing(md_text or "")
    return markdown.markdown(
        spaced,
        # `extra` already includes `tables`; listing `tables` twice used to
        # cause subtle extension-registration races in some markdown versions,
        # so we keep the list minimal: extra + fenced_code + sane_lists.
        extensions=["extra", "fenced_code", "sane_lists"],
        output_format="html5",
    )


def generate_pdf_bytes(title: str, markdown_text: str) -> bytes:
    """Convert a Markdown document to styled PDF bytes via xhtml2pdf (pure Python,
    no native deps — works on Windows out of the box, unlike WeasyPrint)."""
    from xhtml2pdf import pisa

    body_html = _md_to_html(markdown_text or "")
    safe_title = (title or "Document").strip() or "Document"
    full_html = (
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        f"<style>{_PDF_CSS}</style></head><body>"
        f"<h1>{_html_escape(safe_title)}</h1>"
        f"{body_html}"
        "</body></html>"
    )
    buf = io.BytesIO()
    result = pisa.CreatePDF(src=full_html, dest=buf, encoding="utf-8")
    if result.err:
        raise RuntimeError("PDF rendering failed.")
    return buf.getvalue()


def _html_escape(s: str) -> str:
    return (
        s.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def generate_document_pdf(prompt: str) -> tuple[str, bytes]:
    """Pipeline: prompt -> LLM writes Markdown -> we render the PDF. Returns (title, bytes)."""
    # Lazy import so this module doesn't pull rag_service at process start.
    from app.services.rag_service import _chat_complete

    _out = _chat_complete(
        [
            {
                "role": "system",
                "content": (
                    "You are a professional document writer. Given a topic, produce a polished, "
                    "well-structured document in MARKDOWN ONLY.\n\n"
                    "STRUCTURE: Start with a single H1 ('# Title'), then use ## / ### sections, "
                    "bullet lists, numbered lists, and tables where they genuinely help. Be "
                    "substantive and information-dense — avoid filler and waffle.\n\n"
                    "TABLE RULES (CRITICAL — follow exactly):\n"
                    "- Always leave a BLANK LINE before the first table row and after the last row.\n"
                    "- Row format: | Cell | Cell | with pipes at both ends.\n"
                    "- The separator row must use at least 3 dashes per cell: | --- | --- |\n"
                    "- Example of a correctly formatted table:\n\n"
                    "  | Metric | Value |\n"
                    "  | --- | --- |\n"
                    "  | Runs | 1200 |\n"
                    "  | Average | 48.0 |\n\n"
                    "OUTPUT RULES:\n"
                    "- Do NOT wrap the output in code fences.\n"
                    "- Do NOT include any preamble ('Sure!', 'Here is...') or closing remark.\n"
                    "- Output ONLY the Markdown document."
                ),
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.35,
        max_tokens=4096,
    )
    md = (_out or "").strip()

    # Strip accidental ``` fences around the whole doc.
    if md.startswith("```"):
        lines = md.split("\n")
        if lines and lines[0].startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].strip().startswith("```"):
            lines = lines[:-1]
        md = "\n".join(lines)

    # Extract first H1 as title; drop it from body so it isn't duplicated.
    title = "Document"
    body_lines: list[str] = []
    title_found = False
    for line in md.splitlines():
        if not title_found and line.strip().startswith("# "):
            title = line.strip()[2:].strip()
            title_found = True
            continue
        body_lines.append(line)
    body = "\n".join(body_lines).strip() or md

    pdf_bytes = generate_pdf_bytes(title, body)
    return title, pdf_bytes


# ── Excel: LLM → JSON table → openpyxl → .xlsx ──────────────────────────────

def _build_excel(title: str, headers: list, rows: list) -> bytes:
    """Build a styled .xlsx workbook with openpyxl."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = title[:31]  # Excel sheet name max 31 chars

    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill("solid", fgColor="1a1a2e")
    header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    thin = Side(style="thin", color="cccccc")
    bdr = Border(left=thin, right=thin, top=thin, bottom=thin)
    even_fill = PatternFill("solid", fgColor="f2f2f7")

    for col_i, h in enumerate(headers, 1):
        c = ws.cell(row=1, column=col_i, value=h)
        c.font = header_font
        c.fill = header_fill
        c.alignment = header_align
        c.border = bdr

    for row_i, row in enumerate(rows, 2):
        for col_i, val in enumerate(row, 1):
            if col_i > len(headers):
                break
            c = ws.cell(row=row_i, column=col_i, value=val)
            c.alignment = Alignment(vertical="center")
            c.border = bdr
            if row_i % 2 == 0:
                c.fill = even_fill

    for col_i in range(1, len(headers) + 1):
        col_letter = get_column_letter(col_i)
        max_len = len(str(headers[col_i - 1])) if col_i <= len(headers) else 10
        for row_i in range(2, len(rows) + 2):
            v = ws.cell(row=row_i, column=col_i).value
            if v:
                max_len = max(max_len, len(str(v)))
        ws.column_dimensions[col_letter].width = min(max_len + 4, 45)

    ws.freeze_panes = "A2"
    ws.auto_filter.ref = ws.dimensions

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def generate_document_excel(prompt: str) -> tuple[str, bytes]:
    """Pipeline: prompt → LLM writes table JSON → openpyxl → .xlsx bytes."""
    import json as _json
    from app.services.rag_service import _chat_complete

    _out = _chat_complete(
        [
            {
                "role": "system",
                "content": (
                    "You are a data analyst. Given a topic or request, produce a structured "
                    "spreadsheet dataset as JSON ONLY — no prose, no markdown fences.\n\n"
                    "Output EXACTLY this JSON shape:\n"
                    "{\"title\": \"Sheet title\", \"headers\": [\"Col1\", \"Col2\"], "
                    "\"rows\": [[\"val1\", \"val2\"], ...]}\n\n"
                    "Rules:\n"
                    "- title: concise sheet/workbook name\n"
                    "- headers: column names (strings)\n"
                    "- rows: list of lists, one per data row; include 5–30 rows of realistic data\n"
                    "- Consistent data types per column\n"
                    "- No markdown, no fences, no explanations — pure JSON ONLY."
                ),
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.2,
        max_tokens=3000,
    )
    raw = (_out or "").strip()
    if raw.startswith("```"):
        lines = raw.split("\n")
        lines = lines[1:] if lines[0].startswith("```") else lines
        if lines and lines[-1].strip().startswith("```"):
            lines = lines[:-1]
        raw = "\n".join(lines)
    s, e = raw.find("{"), raw.rfind("}")
    if s != -1 and e != -1:
        raw = raw[s : e + 1]
    try:
        data = _json.loads(raw)
    except Exception:
        data = {"title": prompt[:40], "headers": ["Data"], "rows": [["No data generated"]]}

    title = str(data.get("title") or prompt[:40]).strip() or "Sheet"
    headers = [str(h) for h in (data.get("headers") or [])]
    rows = [[str(c) for c in (r or [])] for r in (data.get("rows") or [])]
    if not headers:
        headers = ["Column 1"]
    return title, _build_excel(title, headers, rows)


# ── Word: LLM → Markdown → python-docx → .docx ───────────────────────────────

def _para_with_fmt(para, text: str):
    """Write **bold**, *italic*, `code` inline markup into an existing paragraph."""
    import re as _re
    tokens = _re.split(r"(\*\*(?:[^*]|\*(?!\*))+\*\*|\*[^*]+\*|`[^`]+`)", text)
    for tok in tokens:
        if not tok:
            continue
        if tok.startswith("**") and tok.endswith("**") and len(tok) > 4:
            r = para.add_run(tok[2:-2])
            r.bold = True
        elif tok.startswith("*") and tok.endswith("*") and len(tok) > 2:
            r = para.add_run(tok[1:-1])
            r.italic = True
        elif tok.startswith("`") and tok.endswith("`") and len(tok) > 2:
            r = para.add_run(tok[1:-1])
            r.font.name = "Courier New"
        else:
            para.add_run(tok)


def _add_docx_table(doc, lines: list):
    """Parse markdown table lines and add a formatted Word table."""
    import re as _re

    parsed = []
    for line in lines:
        cells = [c.strip() for c in line.split("|")]
        cells = [c for c in cells if c]
        if not cells:
            continue
        if all(_re.match(r"^[-:]+$", c) for c in cells):
            continue
        parsed.append(cells)
    if not parsed:
        return

    ncols = max(len(r) for r in parsed)
    table = doc.add_table(rows=len(parsed), cols=ncols)
    table.style = "Table Grid"
    for r_idx, row in enumerate(parsed):
        for c_idx, cell_text in enumerate(row):
            if c_idx >= ncols:
                break
            cell = table.rows[r_idx].cells[c_idx]
            para = cell.paragraphs[0]
            _para_with_fmt(para, cell_text)
            if r_idx == 0:
                for run in para.runs:
                    run.bold = True


def _md_to_docx(title: str, body: str) -> bytes:
    """Convert a Markdown body to a styled .docx via python-docx."""
    import io, re as _re
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    for section in doc.sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.25)
        section.right_margin = Inches(1.25)

    doc.add_heading(title, level=0)

    lines = body.split("\n")
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()
        if not stripped:
            i += 1
            continue
        if stripped.startswith("### "):
            doc.add_heading(stripped[4:].strip(), level=3)
            i += 1; continue
        if stripped.startswith("## "):
            doc.add_heading(stripped[3:].strip(), level=2)
            i += 1; continue
        if stripped.startswith("# "):
            doc.add_heading(stripped[2:].strip(), level=1)
            i += 1; continue
        if stripped.startswith("|"):
            tbl_lines = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                tbl_lines.append(lines[i])
                i += 1
            _add_docx_table(doc, tbl_lines)
            doc.add_paragraph()
            continue
        m = _re.match(r"^[-*+]\s+(.*)", stripped)
        if m:
            p = doc.add_paragraph(style="List Bullet")
            _para_with_fmt(p, m.group(1))
            i += 1; continue
        m = _re.match(r"^\d+\.\s+(.*)", stripped)
        if m:
            p = doc.add_paragraph(style="List Number")
            _para_with_fmt(p, m.group(1))
            i += 1; continue
        if _re.match(r"^[-*_]{3,}\s*$", stripped):
            doc.add_paragraph("─" * 50)
            i += 1; continue
        p = doc.add_paragraph()
        _para_with_fmt(p, stripped)
        i += 1

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def generate_document_word(prompt: str) -> tuple[str, bytes]:
    """Pipeline: prompt → LLM writes Markdown → python-docx → .docx bytes."""
    from app.services.rag_service import _chat_complete

    _out = _chat_complete(
        [
            {
                "role": "system",
                "content": (
                    "You are a professional document writer. Given a topic, produce a polished, "
                    "well-structured document in MARKDOWN ONLY.\n\n"
                    "STRUCTURE: Start with a single H1 ('# Title'), then use ## / ### sections, "
                    "bullet lists, numbered lists, and tables where genuinely helpful.\n\n"
                    "TABLE RULES:\n"
                    "- Blank line before and after every table.\n"
                    "- Row format: | Cell | Cell | — separator row: | --- | --- |\n\n"
                    "OUTPUT RULES:\n"
                    "- No code fences, no preamble, no closing remark.\n"
                    "- Output ONLY the Markdown document."
                ),
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.35,
        max_tokens=4096,
    )
    md = (_out or "").strip()
    if md.startswith("```"):
        lines = md.split("\n")
        lines = lines[1:] if lines[0].startswith("```") else lines
        if lines and lines[-1].strip().startswith("```"):
            lines = lines[:-1]
        md = "\n".join(lines)

    title = "Document"
    body_lines: list = []
    title_found = False
    for line in md.splitlines():
        if not title_found and line.strip().startswith("# "):
            title = line.strip()[2:].strip()
            title_found = True
            continue
        body_lines.append(line)
    body = "\n".join(body_lines).strip() or md

    word_bytes = _md_to_docx(title, body)
    return title, word_bytes


# ── PowerPoint: LLM → slide JSON → python-pptx → .pptx ──────────────────────

def _build_ppt(title: str, slides_data: list) -> bytes:
    """Build a clean .pptx from slide JSON using python-pptx."""
    import io
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    try:
        title_slide.placeholders[1].text = "Generated by Close AI"
    except (IndexError, KeyError):
        pass

    # Content slides
    for s in slides_data:
        slide_title = str(s.get("title") or "").strip()
        bullets = [str(b) for b in (s.get("content") or []) if str(b).strip()]

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = slide_title
        try:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for idx, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0
        except (IndexError, KeyError):
            pass

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_document_ppt(prompt: str) -> tuple[str, bytes]:
    """Pipeline: prompt → LLM writes slide JSON → python-pptx → .pptx bytes."""
    import json as _json
    from app.services.rag_service import _chat_complete

    _out = _chat_complete(
        [
            {
                "role": "system",
                "content": (
                    "You are a presentation expert. Given a topic, produce a structured slide deck as JSON ONLY.\n\n"
                    "Output EXACTLY this JSON shape:\n"
                    "{\"title\": \"Deck Title\", \"slides\": ["
                    "{\"title\": \"Slide Title\", \"content\": [\"Bullet 1\", \"Bullet 2\"]}"
                    ", ...]}\n\n"
                    "Rules:\n"
                    "- title: overall deck title\n"
                    "- slides: 6–12 slides; first = intro, last = summary/takeaways\n"
                    "- Each slide: title (short) + content (3–6 bullets, each ≤ 100 chars)\n"
                    "- No markdown, no fences, no explanations — pure JSON ONLY."
                ),
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.3,
        max_tokens=2000,
    )
    raw = (_out or "").strip()
    if raw.startswith("```"):
        lines = raw.split("\n")
        lines = lines[1:] if lines[0].startswith("```") else lines
        if lines and lines[-1].strip().startswith("```"):
            lines = lines[:-1]
        raw = "\n".join(lines)
    s, e = raw.find("{"), raw.rfind("}")
    if s != -1 and e != -1:
        raw = raw[s : e + 1]
    try:
        data = _json.loads(raw)
    except Exception:
        data = {"title": prompt[:40], "slides": [{"title": "Overview", "content": [prompt[:100]]}]}

    title = str(data.get("title") or prompt[:40]).strip() or "Presentation"
    slides = data.get("slides") or []
    return title, _build_ppt(title, slides)
